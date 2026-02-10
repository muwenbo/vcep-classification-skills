#!/usr/bin/env python3
"""
Read PowerPoint presentations (.pptx) and convert to markdown format.

Usage:
    python read_ppt.py <path_to_pptx_file> [--json] [--slide N]

Arguments:
    path_to_pptx_file    Path to the PowerPoint file to read
    --json               Output as JSON instead of markdown (optional)
    --slide N            Read only slide number N (1-indexed, optional)

Output:
    Markdown text with slide structure (default) or JSON structure

Dependencies:
    python-pptx (auto-installed if missing)
"""

import sys
import json
import argparse

def install_dependencies():
    """Install required packages if not present."""
    import subprocess
    try:
        from pptx import Presentation
    except ImportError:
        print("Installing python-pptx...", file=sys.stderr)
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "python-pptx", "-q"],
            stderr=subprocess.DEVNULL
        )

def read_pptx_file(file_path, slide_number=None):
    """
    Read a .pptx file and extract its contents.

    Args:
        file_path: Path to the PowerPoint file
        slide_number: Specific slide to read (1-indexed, None for all)

    Returns:
        Dictionary with presentation structure and content
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        install_dependencies()
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    result = {
        'slide_count': len(prs.slides),
        'slides': []
    }

    for idx, slide in enumerate(prs.slides, 1):
        if slide_number is not None and idx != slide_number:
            continue

        slide_data = {
            'number': idx,
            'title': '',
            'content': [],
            'tables': [],
            'notes': ''
        }

        for shape in slide.shapes:
            # Extract title
            if shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type is not None and ph_type.name in ['TITLE', 'CENTER_TITLE']:
                        if hasattr(shape, 'text'):
                            slide_data['title'] = shape.text.strip()
                            continue
                except Exception:
                    pass

            # Extract text content
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                if text != slide_data['title']:  # Avoid duplicating title
                    slide_data['content'].append(text)

            # Extract tables
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_data['tables'].append(table_data)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text:
                slide_data['notes'] = notes_frame.text.strip()

        result['slides'].append(slide_data)

    return result

def format_as_markdown(data):
    """Convert extracted presentation data to markdown."""
    output = []

    output.append(f"# Presentation ({data['slide_count']} slides)\n")

    for slide in data['slides']:
        output.append(f"\n## Slide {slide['number']}")
        output.append("-" * 40)

        if slide['title']:
            output.append(f"### {slide['title']}\n")

        for content in slide['content']:
            # Format as bullet points if multiple lines
            lines = content.split('\n')
            for line in lines:
                line = line.strip()
                if line:
                    output.append(f"- {line}")
            output.append('')

        # Format tables
        for table in slide['tables']:
            if not table:
                continue
            output.append('')
            # Header row
            header = table[0] if table else []
            output.append('| ' + ' | '.join(header) + ' |')
            output.append('|' + '|'.join(['---'] * len(header)) + '|')
            # Data rows
            for row in table[1:]:
                while len(row) < len(header):
                    row.append('')
                output.append('| ' + ' | '.join(row[:len(header)]) + ' |')
            output.append('')

        # Add notes if present
        if slide['notes']:
            output.append(f"\n> **Speaker Notes:** {slide['notes']}")

        output.append('')

    return '\n'.join(output)

def main():
    parser = argparse.ArgumentParser(
        description='Read PowerPoint files and convert to markdown or JSON'
    )
    parser.add_argument('file_path', help='Path to the PowerPoint file')
    parser.add_argument('--json', action='store_true', help='Output as JSON')
    parser.add_argument('--slide', type=int, help='Read only specific slide (1-indexed)')

    args = parser.parse_args()

    try:
        data = read_pptx_file(args.file_path, args.slide)

        if not data['slides']:
            print("No slides found", file=sys.stderr)
            sys.exit(1)

        if args.json:
            print(json.dumps(data, indent=2, default=str))
        else:
            print(format_as_markdown(data))

    except FileNotFoundError:
        print(f"Error: File not found: {args.file_path}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
