"""Regression tests for the Office reader helpers.

Two defects were fixed in round 15 (2026-08-18):

* read_ppt.py silently dropped every textbox nested inside a group shape.
  ClinGen PVS1 decision trees are built from grouped textboxes, so entire
  flowcharts vanished from the output.
* read_word.py never surfaced images stored under ``word/media/``. Several
  ClinGen supplementary tables ship only as embedded PNGs.

Both tests build their fixtures in-process so they need no external files.
"""

import importlib.util
import unittest
from pathlib import Path

SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"


def _load(name, path):
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


read_ppt = _load("vcep_read_ppt", SCRIPTS / "read_ppt.py")
read_word = _load("vcep_read_word", SCRIPTS / "read_word.py")


class GroupShapeRecursionTests(unittest.TestCase):
    def _fixture(self, tmpdir):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        ).text_frame.text = "TOPLEVEL_TEXT"
        group = slide.shapes.add_group_shape()
        group.shapes.add_textbox(
            Inches(1), Inches(3), Inches(3), Inches(1)
        ).text_frame.text = "GROUPED_TEXT"
        nested = group.shapes.add_group_shape()
        nested.shapes.add_textbox(
            Inches(1), Inches(5), Inches(3), Inches(1)
        ).text_frame.text = "NESTED_GROUP_TEXT"
        path = Path(tmpdir) / "group.pptx"
        prs.save(str(path))
        return path

    def test_text_inside_nested_groups_is_captured(self):
        import tempfile

        with tempfile.TemporaryDirectory() as tmpdir:
            path = self._fixture(tmpdir)
            data = read_ppt.read_pptx_file(str(path))
            content = "\n".join(data["slides"][0]["content"])
            self.assertIn("TOPLEVEL_TEXT", content)
            self.assertIn("GROUPED_TEXT", content)
            self.assertIn("NESTED_GROUP_TEXT", content)


class EmbeddedImageTests(unittest.TestCase):
    # 1x1 transparent PNG
    PNG = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cc\xf8\xcf"
        b"\xc0\xf0\x1f\x00\x05\x05\x02\x00\x9d\xff\xa7\x9f\x00\x00\x00\x00IEND"
        b"\xaeB`\x82"
    )

    def _fixture(self, tmpdir):
        from docx import Document
        from docx.shared import Inches

        png = Path(tmpdir) / "tiny.png"
        png.write_bytes(self.PNG)
        doc = Document()
        doc.add_paragraph("intro")
        doc.add_picture(str(png), width=Inches(1))
        path = Path(tmpdir) / "image.docx"
        doc.save(str(path))
        return path

    def test_embedded_image_is_listed(self):
        import tempfile

        with tempfile.TemporaryDirectory() as tmpdir:
            path = self._fixture(tmpdir)
            data = read_word.read_docx_file(str(path))
            self.assertTrue(data["images"], "embedded image was not surfaced")
            self.assertTrue(data["images"][0]["name"].endswith(".png"))

    def test_extract_media_writes_files(self):
        import tempfile

        with tempfile.TemporaryDirectory() as tmpdir:
            path = self._fixture(tmpdir)
            out = Path(tmpdir) / "media"
            data = read_word.read_docx_file(str(path), str(out))
            saved = data["images"][0]["saved_to"]
            self.assertTrue(Path(saved).is_file())
            self.assertEqual(Path(saved).read_bytes(), self.PNG)


if __name__ == "__main__":
    unittest.main()
